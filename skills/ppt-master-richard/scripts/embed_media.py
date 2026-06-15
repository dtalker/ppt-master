#!/usr/bin/env python3
"""PPT Master · Embed video / audio into an exported PPTX (Step 7.4).

svg_to_pptx.py renders static slides only — it has no path for movies/audio.
This script post-processes the exported .pptx and embeds real, clickable media.

PRIMARY (coordinate-free) mode — named slots
--------------------------------------------
In the SVG, tag the media placeholder rect with ``data-media-slot``:

    <rect data-media-slot="song" x="96" y="236" width="312" height="312" .../>

The slot rect's page (its position among ``sorted(svg_output/*.svg)``) and its
x/y/width/height ARE the embed target — you never type a page number or a
coordinate. Then either:

    python3 embed_media.py <project> --put song=assets/song.mp3 --put demo=assets/demo.mp4

or list slots in ``media.json`` / ``spec_lock.md ## media`` (slot + file only).

Audio is auto-wrapped into a still-poster .mp4 (Keynote won't reliably play
embedded audio); video posters auto-extract; everything scales px→EMU like the
deck (1 canvas-px == slide_width / canvas_width EMU). ffmpeg required for those.

Explicit mode (still supported)
-------------------------------
    {"page": "P02", "file": "...", "x":96, "y":236, "w":312, "h":312, "poster": "..."}
  page = "P<NN>" (1-based position) or an svg basename.

Usage:
    python3 embed_media.py <project>                         # reads media.json / spec_lock
    python3 embed_media.py <project> --put song=foo.mp3 ...  # quick, no manifest
    python3 embed_media.py <project> --pptx in.pptx --out out.pptx
"""
from __future__ import annotations
import argparse, json, re, shutil, subprocess, sys
from pathlib import Path

AUDIO_EXT = {'.mp3', '.wav', '.m4a', '.aac', '.flac', '.ogg', '.oga'}
VIDEO_MIME = {'.mp4': 'video/mp4', '.mov': 'video/quicktime', '.m4v': 'video/mp4',
              '.webm': 'video/webm', '.avi': 'video/x-msvideo'}
CANVAS_DEFAULT_W = 1280


def _have_ffmpeg() -> bool:
    return shutil.which('ffmpeg') is not None


def _canvas_width(project: Path) -> int:
    sl = project / 'spec_lock.md'
    if sl.exists():
        m = re.search(r'viewBox:\s*0\s+0\s+(\d+)\s+\d+', sl.read_text(encoding='utf-8'))
        if m:
            return int(m.group(1))
    return CANVAS_DEFAULT_W


def _svg_files(project: Path) -> list[Path]:
    for d in ('svg_output', 'svg_final'):
        sd = project / d
        if sd.exists():
            return sorted(sd.glob('*.svg'))
    return []


def _slide_order(project: Path) -> list[str]:
    return [p.stem for p in _svg_files(project)]


_ATTR = lambda name, s: (re.search(rf'\b{name}\s*=\s*"([^"]*)"', s) or [None, None])[1]


def _scan_slots(project: Path) -> dict[str, dict]:
    """Find every <... data-media-slot="NAME" x w/width ...> across slides.

    Returns name -> {idx (0-based slide), x, y, w, h}. The first match wins.
    """
    slots: dict[str, dict] = {}
    for idx, svg in enumerate(_svg_files(project)):
        text = svg.read_text(encoding='utf-8')
        for tag in re.findall(r'<[^>]*\bdata-media-slot\s*=\s*"[^"]+"[^>]*>', text):
            name = _ATTR('data-media-slot', tag)
            try:
                x = float(_ATTR('x', tag)); y = float(_ATTR('y', tag))
                w = float(_ATTR('width', tag)); h = float(_ATTR('height', tag))
            except (TypeError, ValueError):
                continue
            if name and name not in slots:
                slots[name] = {'idx': idx, 'x': x, 'y': y, 'w': w, 'h': h, 'svg': svg.stem}
    return slots


def _resolve_index(page: str, order: list[str]) -> int:
    page = str(page).strip()
    m = re.fullmatch(r'[Pp](\d+)', page)
    if m:
        return int(m.group(1)) - 1
    stem = page[:-4] if page.lower().endswith('.svg') else page
    if stem in order:
        return order.index(stem)
    for i, s in enumerate(order):
        if s == stem or s.endswith('_' + stem) or s.endswith(stem):
            return i
    raise KeyError(f'page "{page}" not found among {len(order)} slides')


def _load_manifest(project: Path) -> list[dict]:
    mj = project / 'media.json'
    if mj.exists():
        data = json.loads(mj.read_text(encoding='utf-8'))
        return data if isinstance(data, list) else data.get('media', [])
    sl = project / 'spec_lock.md'
    items: list[dict] = []
    if sl.exists():
        m = re.search(r'^##\s*media\s*$(.*?)(?=^##\s|\Z)', sl.read_text(encoding='utf-8'), re.M | re.S)
        if m:
            for line in m.group(1).splitlines():
                line = line.strip()
                if not line.startswith('-'):
                    continue
                parts = [p.strip() for p in line.lstrip('-').split('|')]
                if len(parts) < 2:
                    continue
                head = parts[0]
                if head.startswith('@'):  # slot form: - @song | file | [poster]
                    d = {'slot': head[1:].strip(), 'file': parts[1]}
                    if len(parts) >= 3 and parts[2]:
                        d['poster'] = parts[2]
                    items.append(d); continue
                if len(parts) < 3:  # page form needs coords
                    continue
                xs = [int(v) for v in re.split(r'[ ,]+', parts[2].strip()) if v]
                if len(xs) != 4:
                    continue
                d = {'page': head, 'file': parts[1], 'x': xs[0], 'y': xs[1], 'w': xs[2], 'h': xs[3]}
                if len(parts) >= 4 and parts[3]:
                    d['poster'] = parts[3]
                items.append(d)
    return items


def _gen_audio_poster(out: Path):
    try:
        from PIL import Image, ImageDraw
    except Exception:
        return None
    W = 640
    im = Image.new('RGB', (W, W), '#0A1F3D'); d = ImageDraw.Draw(im)
    c = W // 2
    d.ellipse([c - 110, c - 110, c + 110, c + 110], outline='#DDB23C', width=6)
    d.ellipse([c - 78, c - 78, c + 78, c + 78], fill='#DDB23C')
    d.polygon([(c - 26, c - 40), (c - 26, c + 40), (c + 44, c)], fill='#0A1F3D')
    im.save(out)
    return out


def _wrap_audio_to_mp4(audio: Path, poster, cache: Path):
    cache.mkdir(parents=True, exist_ok=True)
    if not poster or not Path(poster).exists():
        poster = cache / (audio.stem + '_poster.png'); _gen_audio_poster(poster)
    out = cache / (audio.stem + '_wrapped.mp4')
    if not out.exists():
        subprocess.run(['ffmpeg', '-y', '-loop', '1', '-i', str(poster), '-i', str(audio),
                        '-c:v', 'libx264', '-tune', 'stillimage', '-pix_fmt', 'yuv420p',
                        '-c:a', 'aac', '-b:a', '192k', '-shortest', '-movflags', '+faststart',
                        str(out)], check=True, capture_output=True)
    return out, Path(poster)


def _extract_poster(video: Path, cache: Path):
    if not _have_ffmpeg():
        return None
    cache.mkdir(parents=True, exist_ok=True)
    out = cache / (video.stem + '_poster.jpg')
    if not out.exists():
        try:
            subprocess.run(['ffmpeg', '-y', '-ss', '0.3', '-i', str(video), '-frames:v', '1',
                            '-q:v', '2', str(out)], check=True, capture_output=True)
        except Exception:
            return None
    return out if out.exists() else None


def main():
    ap = argparse.ArgumentParser(description='Embed video/audio into exported PPTX')
    ap.add_argument('project')
    ap.add_argument('--put', action='append', default=[], metavar='slot=file',
                    help='attach a file to a named data-media-slot (repeatable)')
    ap.add_argument('--pptx', help='input pptx (default: newest exports/*.pptx without _media)')
    ap.add_argument('--out', help='output pptx (default: <input>_media.pptx)')
    args = ap.parse_args()

    from pptx import Presentation
    from pptx.util import Emu

    project = Path(args.project).resolve()
    slots = _scan_slots(project)

    manifest = list(_load_manifest(project))
    for spec in args.put:
        if '=' not in spec:
            sys.exit(f'--put expects slot=file, got: {spec}')
        name, f = spec.split('=', 1)
        manifest.append({'slot': name.strip(), 'file': f.strip()})

    if not manifest:
        hint = f' Available slots: {", ".join(slots)}.' if slots else ''
        print(f'[embed_media] No media declared (--put / media.json / spec_lock ## media).{hint}')
        return

    if args.pptx:
        src = Path(args.pptx)
    else:
        cands = [p for p in (project / 'exports').glob('*.pptx') if '_media' not in p.name]
        if not cands:
            sys.exit('[embed_media] No exports/*.pptx found — run svg_to_pptx.py first.')
        src = max(cands, key=lambda p: p.stat().st_mtime)
    out = Path(args.out) if args.out else src.with_name(src.stem + '_media.pptx')

    order = _slide_order(project)
    prs = Presentation(str(src))
    n = len(prs.slides._sldIdLst)
    scale = prs.slide_width / _canvas_width(project)
    cache = project / 'media' / '_wrapped'

    if not _have_ffmpeg():
        print('[embed_media] WARNING: ffmpeg not found — audio cannot be wrapped to video and '
              'video posters cannot be auto-extracted. Provide posters; audio may not play in Keynote.')

    def px(v): return Emu(int(round(v * scale)))

    done = 0
    for item in manifest:
        # --- resolve placement (slot first = no coords needed) ---
        if item.get('slot'):
            s = slots.get(item['slot'])
            if not s:
                print(f"  [skip] slot \"{item['slot']}\": no <rect data-media-slot=\"{item['slot']}\"> "
                      f"in any slide. Known: {', '.join(slots) or '(none)'}"); continue
            idx, x, y, w, h = s['idx'], s['x'], s['y'], s['w'], s['h']
        else:
            try:
                idx = _resolve_index(item['page'], order)
            except Exception as e:
                print(f"  [skip] {item.get('page')}: {e}"); continue
            x, y, w, h = item['x'], item['y'], item['w'], item['h']
        if not (0 <= idx < n):
            print(f"  [skip] {item.get('slot') or item.get('page')}: slide {idx+1} out of range (deck has {n})"); continue

        f = Path(item['file'])
        if not f.is_absolute() and (project / f).exists():
            f = project / f
        if not f.exists():
            print(f"  [skip] {item.get('slot') or item.get('page')}: file not found {item['file']}"); continue

        poster = item.get('poster')
        if poster and not Path(poster).is_absolute() and (project / poster).exists():
            poster = project / poster

        ext = f.suffix.lower()
        if ext in AUDIO_EXT:
            if not _have_ffmpeg():
                print(f"  [skip] audio needs ffmpeg to wrap: {f.name}"); continue
            movie, poster = _wrap_audio_to_mp4(f, poster, cache); mime = 'video/mp4'
        else:
            movie, mime = f, VIDEO_MIME.get(ext, 'video/unknown')
            if not poster or not Path(poster).exists():
                poster = _extract_poster(f, cache)

        kw = dict(mime_type=mime)
        if poster and Path(poster).exists():
            kw['poster_frame_image'] = str(poster)
        prs.slides[idx].shapes.add_movie(str(movie), px(x), px(y), px(w), px(h), **kw)
        tag = item.get('slot') or item.get('page')
        print(f"  [ok] slide {idx+1:>2} ({order[idx] if idx < len(order) else '?'}) "
              f"<- {f.name}  [{tag}]{' (audio→video)' if ext in AUDIO_EXT else ''}")
        done += 1

    if done:
        prs.save(str(out)); print(f"\n[embed_media] Embedded {done} item(s) -> {out}")
    else:
        print('[embed_media] Nothing embedded.')


if __name__ == '__main__':
    main()
